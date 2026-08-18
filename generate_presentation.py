import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Professional IEEE Academic Color Palette
    C_NAVY_DARK = RGBColor(15, 29, 64)       # #0F1D40 - Deep Navy (Header / Primary)
    C_NAVY_BLUE = RGBColor(27, 60, 115)      # #1B3C73 - Sub-header / Accent
    C_BLUE_LIGHT = RGBColor(238, 244, 255)   # #EEF4FF - Soft Card Fill
    C_BLUE_BORDER = RGBColor(190, 212, 245)  # #BED4F5 - Light Accent Border
    C_SLATE_BG = RGBColor(248, 250, 252)     # #F8FAFC - Main Slide Background
    C_WHITE = RGBColor(255, 255, 255)        # #FFFFFF - White Box Fill
    C_TEXT_DARK = RGBColor(30, 41, 59)       # #1E293B - Primary Dark Text
    C_TEXT_MUTED = RGBColor(100, 116, 139)   # #64748B - Subdued/Caption Text
    C_BORDER = RGBColor(226, 232, 240)       # #E2E8F0 - Clean Card Border
    C_ACCENT_BLUE = RGBColor(37, 99, 235)    # #2563EB - Vibrant Accent Blue
    C_ACCENT_TEAL = RGBColor(14, 116, 144)   # #0E7490 - Accent Teal
    C_GREEN = RGBColor(16, 149, 93)          # #10955D - Positive / Low Risk
    C_AMBER = RGBColor(217, 119, 6)          # #D97706 - Moderate Risk / Warning
    C_RED = RGBColor(220, 38, 38)            # #DC2626 - High Risk / Danger
    C_TABLE_HEADER = RGBColor(20, 45, 90)    # Table Header Dark Navy
    C_TABLE_ALT = RGBColor(245, 248, 255)    # Table Alternate Row Light Blue

    FONT_FAMILY = "Calibri"

    def add_base_slide(slide_num, title_text, category_text="FINAL YEAR PROJECT REVIEW — PHASE I"):
        slide = prs.slides.add_slide(blank_layout)
        
        # Base Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_SLATE_BG
        bg.line.fill.background()

        # Top Header Bar
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C_NAVY_DARK
        top_bar.line.color.rgb = C_NAVY_DARK

        # Category / Subtitle
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.0), Inches(0.28))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.name = FONT_FAMILY
        p_c.font.size = Pt(9.5)
        p_c.font.bold = True
        p_c.font.color.rgb = RGBColor(147, 197, 253)

        # Slide Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.0), Inches(0.65))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = FONT_FAMILY
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = C_WHITE

        # Slide Number Badge in Header (Right)
        snum_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.0), Inches(0.32), Inches(0.65), Inches(0.48))
        snum_badge.fill.solid()
        snum_badge.fill.fore_color.rgb = RGBColor(30, 58, 110)
        snum_badge.line.color.rgb = RGBColor(96, 165, 250)
        tf_sn = snum_badge.text_frame
        p_sn = tf_sn.paragraphs[0]
        p_sn.text = f"{slide_num:02d}"
        p_sn.font.name = FONT_FAMILY
        p_sn.font.size = Pt(13)
        p_sn.font.bold = True
        p_sn.font.color.rgb = C_WHITE
        p_sn.alignment = PP_ALIGN.CENTER

        # Accent Line below header
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.12), Inches(13.333), Inches(0.04))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = C_ACCENT_BLUE
        accent_line.line.color.rgb = C_ACCENT_BLUE

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.18), Inches(11.733), Inches(0.25))
        tf_f = footer_box.text_frame
        tf_f.word_wrap = True
        tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0
        p_f = tf_f.paragraphs[0]
        p_f.text = "Department of Information Technology  |  Confidence-Driven Probabilistic ML & RAG Energy Planning Framework"
        p_f.font.name = FONT_FAMILY
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = C_TEXT_MUTED

        return slide

    def add_card(slide, left, top, width, height, title=None, bg_color=C_WHITE, border_color=C_BORDER, title_color=C_NAVY_DARK, title_size=12):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)

        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.35))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = FONT_FAMILY
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = title_color
        return card

    def style_table(table, header_bg=C_TABLE_HEADER, alt_bg=C_TABLE_ALT):
        for col_idx, col in enumerate(table.columns):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_FAMILY
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
        
        for row_idx in range(1, len(table.rows)):
            row_bg = alt_bg if row_idx % 2 == 1 else C_WHITE
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                for p in cell.text_frame.paragraphs:
                    p.font.name = FONT_FAMILY
                    p.font.size = Pt(9)
                    p.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(245, 248, 253)
    bg1.line.fill.background()

    top_header = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.6))
    top_header.fill.solid()
    top_header.fill.fore_color.rgb = C_NAVY_DARK
    top_header.line.color.rgb = C_NAVY_DARK

    sub_tb = slide1.shapes.add_textbox(Inches(0.9), Inches(0.35), Inches(11.533), Inches(0.35))
    tf_sub = sub_tb.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "FINAL YEAR PROJECT REVIEW — PHASE I (ACADEMIC YEAR 2026–2027)"
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(12)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(147, 197, 253)

    title_tb = slide1.shapes.add_textbox(Inches(0.9), Inches(0.75), Inches(11.533), Inches(1.5))
    tf_main = title_tb.text_frame
    tf_main.word_wrap = True
    p_main = tf_main.paragraphs[0]
    p_main.text = "Confidence-Driven Probabilistic Machine Learning Framework for Electricity Demand–Supply Gap Forecasting and RAG-Based Energy Planning"
    p_main.font.name = FONT_FAMILY
    p_main.font.size = Pt(21)
    p_main.font.bold = True
    p_main.font.color.rgb = C_WHITE

    # Ribbon Tag
    ribbon = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.8), Inches(11.533), Inches(0.42))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = C_BLUE_LIGHT
    ribbon.line.color.rgb = C_BLUE_BORDER
    tf_r = ribbon.text_frame
    p_r = tf_r.paragraphs[0]
    p_r.text = "Domain: Machine Learning & Deep Learning  •  Probabilistic Uncertainty Quantification  •  RAG Decision Support"
    p_r.font.name = FONT_FAMILY
    p_r.font.size = Pt(10.5)
    p_r.font.bold = True
    p_r.font.color.rgb = C_NAVY_BLUE
    p_r.alignment = PP_ALIGN.CENTER

    # Presenters Card
    add_card(slide1, Inches(0.9), Inches(3.45), Inches(5.6), Inches(3.5), title="STUDENT PRESENTERS", bg_color=C_WHITE, border_color=C_BORDER, title_size=13)
    stud_tb = slide1.shapes.add_textbox(Inches(1.15), Inches(4.0), Inches(5.1), Inches(2.8))
    tf_stud = stud_tb.text_frame
    tf_stud.word_wrap = True
    
    students = [
        ("S. Udayakumar", "Register Number: 43120227", "Model Architecture & Time-Series Engineering"),
        ("D. Tamilarasu", "Register Number: 43120224", "Uncertainty Estimation & Gap Analytics"),
        ("Suganth Kesavan", "Register Number: 43120216", "RAG Pipeline & Decision Support System"),
    ]
    for idx, (name, regno, role) in enumerate(students):
        p1 = tf_stud.add_paragraph() if idx > 0 else tf_stud.paragraphs[0]
        p1.text = f"•  {name}  ({regno.split(': ')[1]})"
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY_DARK
        p1.space_before = Pt(6) if idx > 0 else Pt(0)
        
        p2 = tf_stud.add_paragraph()
        p2.text = f"    Role: {role}"
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_TEXT_MUTED

    # Guide & Department Card
    add_card(slide1, Inches(6.833), Inches(3.45), Inches(5.6), Inches(3.5), title="PROJECT SUPERVISOR & INSTITUTION", bg_color=C_WHITE, border_color=C_BORDER, title_size=13)
    guide_tb = slide1.shapes.add_textbox(Inches(7.1), Inches(4.0), Inches(5.1), Inches(2.8))
    tf_guide = guide_tb.text_frame
    tf_guide.word_wrap = True

    guide_items = [
        ("Project Supervisor:", "Dr. Kamatchi K. S"),
        ("Academic Designation:", "Associate Professor"),
        ("Department:", "Department of Information Technology"),
        ("Focus Area:", "Energy Analytics, Deep Learning & Applied AI"),
        ("Academic Year:", "2026 – 2027 (Final Year Engineering Project)"),
    ]
    for idx, (label, val) in enumerate(guide_items):
        p1 = tf_guide.add_paragraph() if idx > 0 else tf_guide.paragraphs[0]
        p1.text = f"{label} "
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(11.5)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY_BLUE
        p1.space_before = Pt(6) if idx > 0 else Pt(0)
        
        run = p1.add_run()
        run.text = val
        run.font.name = FONT_FAMILY
        run.font.size = Pt(11.5)
        run.font.bold = (label == "Project Supervisor:")
        run.font.color.rgb = C_NAVY_DARK if label == "Project Supervisor:" else C_TEXT_DARK

    # =========================================================================
    # SLIDE 2: AGENDA
    # =========================================================================
    slide2 = add_base_slide(2, "Presentation Agenda & Review Structure")
    
    agenda_topics = [
        ("01", "Introduction & Background", "Electricity demand growth, renewable intermittency, and the balancing dilemma."),
        ("02", "Problem Statement & Gaps", "Identified limitations in isolated forecasting and unquantified prediction uncertainty."),
        ("03", "Aim, Objectives & Scope", "Project goals, 7 specific technical objectives, and Tamil Nadu grid scope."),
        ("04", "Expected Technical Outcomes", "Dual forecasting, 90/95/99% PIs, gap quantification, and RAG energy planning."),
        ("05", "Literature Survey", "Analysis of 5 major benchmark papers with methodology and limitation mapping."),
        ("06", "Research Gap Analysis", "Structured comparative matrix contrasting conventional approaches with our framework."),
        ("07", "Feasibility & Risk Assessment", "Technical, economic, operational, and environmental feasibility with risk matrix."),
        ("08", "System Architecture & Flow", "Decoupled dual-stream forecasting architecture and RAG knowledge retrieval pipeline."),
        ("09", "Materials, Methods & Data Analysis", "Datasets, feature engineering, ML algorithms, and statistical uncertainty formulation."),
        ("10", "Validation, Timeline & Future Work", "Jan–Mar 2026 test benchmark, model performance metrics, Gantt chart, and Phase II plan."),
    ]

    for idx, (num, title, desc) in enumerate(agenda_topics):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.35 + row * 1.1)
        
        add_card(slide2, left, top, Inches(5.75), Inches(0.98), bg_color=C_WHITE, border_color=C_BORDER)
        
        badge = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.12), top + Inches(0.12), Inches(0.72), Inches(0.72))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_BLUE_LIGHT
        badge.line.color.rgb = C_BLUE_BORDER
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = num
        p_b.font.name = FONT_FAMILY
        p_b.font.size = Pt(13)
        p_b.font.bold = True
        p_b.font.color.rgb = C_NAVY_BLUE
        p_b.alignment = PP_ALIGN.CENTER
        
        tb = slide2.shapes.add_textbox(left + Inches(0.95), top + Inches(0.08), Inches(4.65), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY_DARK
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 3: INTRODUCTION & BACKGROUND
    # =========================================================================
    slide3 = add_base_slide(3, "Introduction & Background: Power Grid Dynamics")

    intro_cards = [
        ("1. Rising Electricity Demand & Urbanization", [
            "Rapid economic expansion, EV adoption, industrial automation, and digital infrastructure drive steep surges in power demand.",
            "Peak load swings exhibit sharp seasonal variations and extreme non-linear intraday fluctuations."
        ]),
        ("2. Renewable Integration & Intermittency", [
            "Accelerated transition towards carbon neutrality via massive integration of Solar PV and Wind energy systems.",
            "Generation is strictly weather-dependent (solar irradiance, wind velocity, ambient temperature, monsoon precipitation)."
        ]),
        ("3. The Demand–Supply Balancing Dilemma", [
            "Severe operational challenge in maintaining instantaneous generation-load equilibrium at State Load Despatch Centres (SLDC).",
            "Unsynchronized demand surges and generation drop-offs cause acute grid frequency deviations and reserve strain."
        ]),
        ("4. Paradigm Shift: Uncertainty & RAG Planning", [
            "Conventional point forecasts fail to provide operational confidence margins or risk boundaries for spinning reserves.",
            "Critical need for probabilistic prediction intervals (90%, 95%, 99%) coupled with knowledge-grounded RAG energy planning."
        ])
    ]

    for idx, (title, points) in enumerate(intro_cards):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.35 + row * 2.5)
        
        add_card(slide3, left, top, Inches(5.75), Inches(2.35), title=title, bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
        
        tb = slide3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.52), Inches(5.35), Inches(1.75))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for p_idx, pt in enumerate(points):
            p = tf.add_paragraph() if p_idx > 0 else tf.paragraphs[0]
            p.text = f"•  {pt}"
            p.font.name = FONT_FAMILY
            p.font.size = Pt(10)
            p.font.color.rgb = C_TEXT_DARK
            p.space_before = Pt(5) if p_idx > 0 else Pt(0)

    bot_card = add_card(slide3, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.55), bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER)
    tb_b = slide3.shapes.add_textbox(Inches(0.9), Inches(6.52), Inches(11.5), Inches(0.4))
    tf_b = tb_b.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "Key Insight: Safe grid planning requires simultaneous demand-supply gap forecasting, multi-level confidence bounds, and actionable policy guidance."
    p_b.font.name = FONT_FAMILY
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = C_NAVY_BLUE
    p_b.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 4: PROBLEM STATEMENT
    # =========================================================================
    slide4 = add_base_slide(4, "Problem Statement: Key Limitations in Existing Systems")

    problems = [
        ("Problem 1: Isolated & Siloed Forecasting", 
         "Existing literature overwhelmingly treats electricity demand and renewable supply as separate problems.",
         "Failing to model both dimensions simultaneously prevents grid planners from knowing the net imbalance in advance."),
        
        ("Problem 2: Deterministic Point Prediction Deficits", 
         "Standard ML/DL systems produce single deterministic values (point forecasts) without confidence bounds.",
         "Point predictions fail to communicate forecast uncertainty, leaving operators blind to tail-risk extreme scenarios."),
        
        ("Problem 3: Absence of Integrated Gap Forecasting", 
         "Direct mathematical gap forecasting with combined uncertainty quantification is practically non-existent in state grids.",
         "Without a synchronized gap model, reserve capacity allocation is either over-provisioned (wasteful) or under-provisioned (outages)."),
        
        ("Problem 4: Disconnect from Actionable Energy Planning", 
         "Existing forecasting algorithms output raw numbers and terminate without contextualized operational guidance.",
         "Grid operators lack automated, policy-grounded decision support to translate numerical shortages into concrete actions.")
    ]

    for idx, (title, desc1, desc2) in enumerate(problems):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.35 + row * 2.25)
        
        add_card(slide4, left, top, Inches(5.75), Inches(2.1), title=title, bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
        
        tb = slide4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), Inches(5.35), Inches(1.45))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p1 = tf.paragraphs[0]
        p1.text = f"•  {desc1}"
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(10)
        p1.font.color.rgb = C_TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = f"•  {desc2}"
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_TEXT_DARK
        p2.space_before = Pt(4)

    core_card = add_card(slide4, Inches(0.8), Inches(5.95), Inches(11.7), Inches(1.05), bg_color=RGBColor(254, 242, 242), border_color=RGBColor(254, 202, 202))
    tb_c = slide4.shapes.add_textbox(Inches(1.0), Inches(6.02), Inches(11.3), Inches(0.9))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "CORE RESEARCH OBJECTIVE ADDRESSED:"
    p_c1.font.name = FONT_FAMILY
    p_c1.font.size = Pt(10.5)
    p_c1.font.bold = True
    p_c1.font.color.rgb = C_RED
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "To bridge this gap by designing an end-to-end framework integrating: (1) Decoupled Demand and Supply ML Forecasting, (2) Probabilistic Uncertainty & Prediction Intervals (90%, 95%, 99%), (3) Demand–Supply Gap Quantification, and (4) RAG-based Energy Planning Recommendations."
    p_c2.font.name = FONT_FAMILY
    p_c2.font.size = Pt(10)
    p_c2.font.color.rgb = C_TEXT_DARK
    p_c2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 5: AIM, OBJECTIVES & SCOPE
    # =========================================================================
    slide5 = add_base_slide(5, "Project Aim, Key Objectives & Operational Scope")

    add_card(slide5, Inches(0.8), Inches(1.35), Inches(5.4), Inches(2.55), title="PROJECT AIM", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    tb_aim = slide5.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.0), Inches(1.95))
    tf_aim = tb_aim.text_frame
    tf_aim.word_wrap = True
    p_aim = tf_aim.paragraphs[0]
    p_aim.text = "To develop a confidence-driven probabilistic machine learning framework that forecasts electricity demand and source-wise supply, quantifies prediction uncertainty through multi-level prediction intervals, estimates the future demand–supply gap, and provides knowledge-grounded energy planning recommendations via Retrieval-Augmented Generation (RAG)."
    p_aim.font.name = FONT_FAMILY
    p_aim.font.size = Pt(10)
    p_aim.font.color.rgb = C_TEXT_DARK

    add_card(slide5, Inches(0.8), Inches(4.05), Inches(5.4), Inches(2.95), title="PROJECT SCOPE & BOUNDARIES", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    tb_scp = slide5.shapes.add_textbox(Inches(1.0), Inches(4.55), Inches(5.0), Inches(2.35))
    tf_scp = tb_scp.text_frame
    tf_scp.word_wrap = True
    scope_pts = [
        ("Geographical Focus:", "Tamil Nadu Power System & State Load Despatch Centre (SLDC)."),
        ("Temporal Scope:", "129 monthly observations (April 2015 – December 2025)."),
        ("Forecast Horizon:", "3-month forward operational planning window (Jan–Mar 2026)."),
        ("Supply Sources (8):", "Coal, Oil & Gas, Nuclear, Hydro, Solar, Wind, Small Hydro, Bio-Power."),
        ("Meteorological Data:", "Temperature, Relative Humidity, Rainfall, Solar Irradiance."),
    ]
    for s_idx, (lbl, val) in enumerate(scope_pts):
        p = tf_scp.add_paragraph() if s_idx > 0 else tf_scp.paragraphs[0]
        p.text = f"•  {lbl} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_BLUE
        p.space_before = Pt(3) if s_idx > 0 else Pt(0)
        run = p.add_run()
        run.text = val
        run.font.name = FONT_FAMILY
        run.font.size = Pt(9.5)
        run.font.bold = False
        run.font.color.rgb = C_TEXT_DARK

    add_card(slide5, Inches(6.4), Inches(1.35), Inches(6.1), Inches(5.65), title="SPECIFIC TECHNICAL OBJECTIVES", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    tb_obj = slide5.shapes.add_textbox(Inches(6.6), Inches(1.85), Inches(5.7), Inches(5.0))
    tf_obj = tb_obj.text_frame
    tf_obj.word_wrap = True

    objs = [
        ("1. Demand Forecasting:", "Develop ML/DL models (RF, XGBoost, LightGBM, LSTM) to forecast monthly electricity demand using lag and weather features."),
        ("2. Supply Forecasting:", "Train independent ML/DL models to forecast monthly source-wise power generation across thermal, nuclear, and renewable sources."),
        ("3. Uncertainty Quantification:", "Estimate prediction uncertainty using the historical validation residual distribution (standard error sigma)."),
        ("4. Multi-Level Prediction Intervals:", "Construct 90%, 95%, and 99% symmetric confidence bounds around point forecasts using standard normal critical values."),
        ("5. Demand–Supply Gap Analysis:", "Calculate the net operational gap (Gap = Predicted Demand - Predicted Supply) and combine standard errors (sigma_G = sqrt(sigma_D^2 + sigma_S^2))."),
        ("6. Risk & Condition Classification:", "Perform deterministic threshold classification (Shortage/Surplus/Balanced; Low <3000 MU, Moderate 3000-4500 MU, High >4500 MU)."),
        ("7. RAG Knowledge Planning:", "Retrieve trusted energy policies, grid regulations, and CEA guidelines from a FAISS vector database to generate LLM-powered planning advice."),
    ]
    for o_idx, (otitle, odesc) in enumerate(objs):
        p = tf_obj.add_paragraph() if o_idx > 0 else tf_obj.paragraphs[0]
        p.text = f"{otitle} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_DARK
        p.space_before = Pt(4) if o_idx > 0 else Pt(0)
        run = p.add_run()
        run.text = odesc
        run.font.name = FONT_FAMILY
        run.font.size = Pt(9.5)
        run.font.bold = False
        run.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 6: EXPECTED OUTCOMES
    # =========================================================================
    slide6 = add_base_slide(6, "Expected Technical Outcomes & System Deliverables")

    outcomes = [
        ("1. Accurate Demand Forecasting Engine",
         "High-accuracy monthly demand point forecasts.",
         "LSTM outperforms baselines with R² = 0.78, MAE = 368.72 MU, MAPE = 3.41%."),
        
        ("2. Source-Wise Supply Forecasting Engine",
         "Independent forecasting of 8 energy generation sources.",
         "LSTM captures renewable intermittency with R² = 0.56, MAE = 731.56 MU, MAPE = 6.30%."),
        
        ("3. Multi-Level Prediction Intervals",
         "Quantified uncertainty with 90%, 95%, and 99% confidence bands.",
         "Provides grid operators with rigorous statistical boundaries for tail-risk planning."),
        
        ("4. Mathematical Demand–Supply Gap Estimation",
         "Direct quantification of net state electricity deficit/surplus.",
         "Propagated standard error: sigma_G = sqrt(sigma_D^2 + sigma_S^2) = 1,070.88 MU."),
        
        ("5. Deterministic Risk & Threshold Classification",
         "Clear operational condition status: Shortage, Balanced, or Surplus.",
         "Threshold severity bands: Low (<3000 MU), Moderate (3000–4500 MU), High (>4500 MU)."),
        
        ("6. RAG-Powered Energy Planning Advice",
         "Knowledge-grounded recommendations synthesized by LLM.",
         "Grounded in CEA, TNERC regulations, spinning reserves, and demand response protocols.")
    ]

    for idx, (title, pt1, pt2) in enumerate(outcomes):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.35 + row * 1.85)
        
        add_card(slide6, left, top, Inches(5.75), Inches(1.7), title=title, bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
        
        tb = slide6.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), Inches(5.35), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p1 = tf.paragraphs[0]
        p1.text = f"•  {pt1}"
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(9.5)
        p1.font.color.rgb = C_TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = f"•  {pt2}"
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = C_TEXT_DARK
        p2.space_before = Pt(3)

    # =========================================================================
    # SLIDE 7: LITERATURE SURVEY
    # =========================================================================
    slide7 = add_base_slide(7, "Literature Survey: Summary of Key Benchmark Studies")

    rows = 6
    cols = 5
    top_t = Inches(1.35)
    left_t = Inches(0.8)
    width_t = Inches(11.733)
    height_t = Inches(5.5)

    table_shape = slide7.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
    table = table_shape.table
    
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.6)
    table.columns[4].width = Inches(2.533)

    headers = ["Author / Year", "Method / Model", "Application Domain", "Major Finding", "Identified Limitation"]
    for idx, h in enumerate(headers):
        table.cell(0, idx).text = h

    lit_data = [
        ("Chisale et al.\n(2025) [5]", "ECEEMDAN-BiGRU &\nQuantile Regression", "Electricity Demand\nForecasting", "Signal decomposition effectively extracts non-linear trends; quantile regression yields robust demand intervals.", "Focuses exclusively on demand; completely omits supply-side forecasting & gap estimation."),
        ("Cao et al.\n(2022) [15]", "Transformer-Guided\nState Space (TPEDF)", "Probabilistic Demand\nForecasting", "Transformer attention captures long-term dependencies while state-space model quantifies uncertainty distributions.", "Single-side forecasting; lacks demand-supply integration and actionable decision support."),
        ("Terrón-Serrano et al.\n(2026) [13]", "Shared Probabilistic\nDistributions", "Day-Ahead Demand &\nRenewables", "Joint probabilistic modeling of load and renewables using weather data improves reserve scheduling efficiency.", "Does not compute operational gap intervals or incorporate RAG-based regulatory planning."),
        ("Ni et al.\n(2024) [16]", "ChatGrid:\nLLM + RAG Pipeline", "Power Dispatching\nKnowledge Q&A", "Retrieval-Augmented Generation grounds LLM answers in grid dispatch rules, dramatically cutting hallucinations.", "Standalone conversational Q&A system; lacks integration with numerical ML forecasting engines."),
        ("Debnath et al.\n(2025) [29]", "Weather-Informed\nLSTM & Transformer", "Extreme Weather Grid\nLoad Forecasting", "Incorporating multi-source weather variables (temperature, humidity) significantly reduces peak load forecast errors.", "Deterministic point predictions only; no confidence intervals or supply-gap risk assessment.")
    ]

    for row_idx, data_row in enumerate(lit_data):
        for col_idx, text_val in enumerate(data_row):
            table.cell(row_idx + 1, col_idx).text = text_val

    style_table(table)

    # =========================================================================
    # SLIDE 8: RESEARCH GAP ANALYSIS
    # =========================================================================
    slide8 = add_base_slide(8, "Research Gap Analysis: Conventional vs Proposed Framework")

    rows = 7
    cols = 4
    top_t = Inches(1.35)
    left_t = Inches(0.8)
    width_t = Inches(11.733)
    height_t = Inches(4.7)

    table_shape8 = slide8.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
    table8 = table_shape8.table
    
    table8.columns[0].width = Inches(2.2)
    table8.columns[1].width = Inches(3.1)
    table8.columns[2].width = Inches(3.2)
    table8.columns[3].width = Inches(3.233)

    headers8 = ["Evaluation Dimension", "Existing Approaches", "Proposed Framework", "Research Gap Addressed"]
    for idx, h in enumerate(headers8):
        table8.cell(0, idx).text = h

    gap_data = [
        ("Forecasting Scope", "Forecasts demand or supply in isolation (siloed models)", "Decoupled, synchronized dual-stream forecasting (Demand + Supply)", "Eliminates operational mismatch between generation and load"),
        ("Prediction Nature", "Deterministic point forecasts (single values)", "Confidence-driven probabilistic forecasting with multi-level bounds", "Quantifies prediction uncertainty for tail-risk management"),
        ("Prediction Intervals", "Rarely implemented, or single-level (e.g., 95% only)", "Multi-level 90%, 95%, and 99% symmetric prediction intervals", "Offers customizable risk tolerance margins for SLDC operators"),
        ("Demand–Supply Gap", "Heuristic or post-hoc comparison without uncertainty", "Dedicated mathematical gap engine with combined sigma_G propagation", "Provides rigorous statistical confidence bounds for net deficit/surplus"),
        ("Risk Assessment", "Manual, ad-hoc threshold inspection by engineers", "Deterministic 3-tier severity classification (Low, Moderate, High)", "Standardizes automated dispatch alerts and operational priority"),
        ("Energy Planning", "Disconnected static rulebooks or ungrounded LLMs", "Integrated RAG engine with FAISS vector retrieval & grounded LLM", "Bridges numerical forecast outputs directly to policy-grounded actions")
    ]

    for row_idx, data_row in enumerate(gap_data):
        for col_idx, text_val in enumerate(data_row):
            table8.cell(row_idx + 1, col_idx).text = text_val

    style_table(table8)

    gap_box = add_card(slide8, Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.85), bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER)
    tb_gb = slide8.shapes.add_textbox(Inches(0.95), Inches(6.25), Inches(11.4), Inches(0.75))
    tf_gb = tb_gb.text_frame
    tf_gb.word_wrap = True
    p_gb1 = tf_gb.paragraphs[0]
    p_gb1.text = "Key Innovation Summary:"
    p_gb1.font.name = FONT_FAMILY
    p_gb1.font.size = Pt(10.5)
    p_gb1.font.bold = True
    p_gb1.font.color.rgb = C_NAVY_BLUE
    
    p_gb2 = tf_gb.add_paragraph()
    p_gb2.text = "Our framework uniquely closes the loop: Data Ingestion -> Decoupled ML Forecasting -> Multi-Level Residual Uncertainty -> Mathematical Gap Estimation -> Deterministic Risk Classification -> RAG-Grounded Energy Planning."
    p_gb2.font.name = FONT_FAMILY
    p_gb2.font.size = Pt(9.5)
    p_gb2.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 9: FEASIBILITY ANALYSIS
    # =========================================================================
    slide9 = add_base_slide(9, "Feasibility Analysis: Technical, Economic, Operational & Environmental")

    feas_cards = [
        ("Technical Feasibility", [
            "Leverages proven deep learning architectures (LSTM) and gradient boosting frameworks (XGBoost, LightGBM, Random Forest).",
            "Residual uncertainty modeling uses robust statistical principles (standard normal critical values z_0.95, z_0.975, z_0.995).",
            "RAG pipeline is built using efficient FAISS vector indexing and state-of-the-art quantized LLMs (Llama-3/Mistral/Ollama)."
        ]),
        ("Economic Feasibility", [
            "Zero data acquisition costs: utilizes openly accessible public datasets from CEA, Grid-India, and NASA POWER.",
            "100% open-source software stack (Python, TensorFlow, Scikit-learn, FAISS, Streamlit) eliminates proprietary licensing fees.",
            "Minimal computational footprint: training and inference execute seamlessly on standard workstation GPUs / multi-core CPUs."
        ]),
        ("Operational Feasibility", [
            "Seamlessly integrates into existing SLDC and utility operational workflows via intuitive dashboard visualizations.",
            "Translates complex probabilistic interval math into clear, color-coded risk alerts and transparent numerical bounds.",
            "Maintains human-in-the-loop governance: RAG provides actionable recommendations while engineers retain final dispatch control."
        ]),
        ("Environmental Feasibility", [
            "Facilitates maximum integration of intermittent renewable sources (solar and wind) into the state energy mix.",
            "Minimizes fossil-fuel spinning reserve wastage and renewable curtailment through accurate 3-month lookahead gap forecasts.",
            "Directly aligns with India's carbon neutrality targets and state-level renewable purchase obligations (RPO)."
        ])
    ]

    for idx, (title, points) in enumerate(feas_cards):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.35 + row * 2.7)
        
        add_card(slide9, left, top, Inches(5.75), Inches(2.55), title=title, bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
        
        tb = slide9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), Inches(5.35), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for p_idx, pt in enumerate(points):
            p = tf.add_paragraph() if p_idx > 0 else tf.paragraphs[0]
            p.text = f"•  {pt}"
            p.font.name = FONT_FAMILY
            p.font.size = Pt(9.5)
            p.font.color.rgb = C_TEXT_DARK
            p.space_before = Pt(5) if p_idx > 0 else Pt(0)

    # =========================================================================
    # SLIDE 10: RESOURCE REQUIREMENTS & RISK ASSESSMENT
    # =========================================================================
    slide10 = add_base_slide(10, "Resource Requirements & Risk Assessment Matrix")

    add_card(slide10, Inches(0.8), Inches(1.35), Inches(5.4), Inches(5.65), title="SYSTEM RESOURCE REQUIREMENTS", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    tb_res = slide10.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.0), Inches(5.0))
    tf_res = tb_res.text_frame
    tf_res.word_wrap = True

    res_items = [
        ("Hardware Requirements:", [
            "Processor: Multi-core CPU (Intel Core i5/i7/i9 or AMD Ryzen 5/7/9)",
            "System Memory: 16 GB DDR4/DDR5 RAM (minimum)",
            "Graphics Accelerator: NVIDIA GPU (CUDA-compatible) for deep learning and vector operations",
            "Storage: 50 GB SSD for datasets, trained model weights, and FAISS indices"
        ]),
        ("Software & Development Stack:", [
            "Operating System: Windows 11 / Linux (Ubuntu 22.04 LTS)",
            "Programming Runtime: Python 3.11 / Python 3.12",
            "ML/DL Libraries: TensorFlow/Keras, Scikit-learn, XGBoost, LightGBM",
            "Data & Scientific: Pandas, NumPy, SciPy, Joblib",
            "RAG & Vector Database: FAISS (CPU/GPU), LangChain / LlamaIndex, Ollama / HuggingFace Transformers",
            "Development Environment: Jupyter Notebook, VS Code / Antigravity IDE"
        ])
    ]

    for s_idx, (sec_title, sec_pts) in enumerate(res_items):
        p_sec = tf_res.add_paragraph() if s_idx > 0 else tf_res.paragraphs[0]
        p_sec.text = sec_title
        p_sec.font.name = FONT_FAMILY
        p_sec.font.size = Pt(10.5)
        p_sec.font.bold = True
        p_sec.font.color.rgb = C_NAVY_BLUE
        p_sec.space_before = Pt(8) if s_idx > 0 else Pt(0)
        
        for pt in sec_pts:
            p_pt = tf_res.add_paragraph()
            p_pt.text = f"•  {pt}"
            p_pt.font.name = FONT_FAMILY
            p_pt.font.size = Pt(9)
            p_pt.font.color.rgb = C_TEXT_DARK
            p_pt.space_before = Pt(2)

    add_card(slide10, Inches(6.4), Inches(1.35), Inches(6.1), Inches(5.65), title="RISK ASSESSMENT & MITIGATION MATRIX", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    
    rows_r = 5
    cols_r = 3
    t_shape_r = slide10.shapes.add_table(rows_r, cols_r, Inches(6.55), Inches(1.9), Inches(5.8), Inches(4.9))
    tbl_r = t_shape_r.table
    tbl_r.columns[0].width = Inches(1.8)
    tbl_r.columns[1].width = Inches(0.9)
    tbl_r.columns[2].width = Inches(3.1)

    r_headers = ["Identified Risk", "Impact", "Mitigation Strategy"]
    for idx, h in enumerate(r_headers):
        tbl_r.cell(0, idx).text = h

    risk_data = [
        ("High Renewable Weather Volatility", "High", "Integrate multi-source weather data (temperature, irradiance, humidity) & multi-tier 90/95/99% prediction intervals."),
        ("Overfitting on Seasonal Peaks", "Medium", "Implement time-series chronological train/test split, lag features, rolling windows, and cyclical harmonic encoders."),
        ("LLM Hallucination in Planning", "High", "Strict RAG architecture with top-k vector retrieval from authoritative energy policies, temperature clamping & prompt constraints."),
        ("Data Inconsistency across Sources", "Medium", "Robust preprocessing pipeline with automated date normalization, missing-value interpolation, and unit harmonization.")
    ]

    for row_idx, r_row in enumerate(risk_data):
        for col_idx, r_val in enumerate(r_row):
            tbl_r.cell(row_idx + 1, col_idx).text = r_val

    style_table(tbl_r)

    # =========================================================================
    # SLIDE 11: PROPOSED SYSTEM OVERVIEW (ARCHITECTURE DIAGRAM)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    bg11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg11.fill.solid()
    bg11.fill.fore_color.rgb = C_SLATE_BG
    bg11.line.fill.background()

    top_bar11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.12))
    top_bar11.fill.solid()
    top_bar11.fill.fore_color.rgb = C_NAVY_DARK
    top_bar11.line.color.rgb = C_NAVY_DARK

    cat_box11 = slide11.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.0), Inches(0.28))
    tf_c11 = cat_box11.text_frame
    p_c11 = tf_c11.paragraphs[0]
    p_c11.text = "FINAL YEAR PROJECT REVIEW — SYSTEM ARCHITECTURE"
    p_c11.font.name = FONT_FAMILY
    p_c11.font.size = Pt(9.5)
    p_c11.font.bold = True
    p_c11.font.color.rgb = RGBColor(147, 197, 253)

    title_box11 = slide11.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.0), Inches(0.65))
    tf_t11 = title_box11.text_frame
    p_t11 = tf_t11.paragraphs[0]
    p_t11.text = "Proposed System Overview & End-to-End Architecture"
    p_t11.font.name = FONT_FAMILY
    p_t11.font.size = Pt(20)
    p_t11.font.bold = True
    p_t11.font.color.rgb = C_WHITE

    accent_line11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.12), Inches(13.333), Inches(0.04))
    accent_line11.fill.solid()
    accent_line11.fill.fore_color.rgb = C_ACCENT_BLUE
    accent_line11.line.color.rgb = C_ACCENT_BLUE

    # Architecture Diagram Canvas (Visual Blocks & Strict Decoupled Structure)
    # 1. Data Ingestion Layer (Top)
    d_box = add_card(slide11, Inches(2.5), Inches(1.3), Inches(8.333), Inches(0.65), bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER)
    tb_d = slide11.shapes.add_textbox(Inches(2.6), Inches(1.35), Inches(8.133), Inches(0.55))
    tf_d = tb_d.text_frame
    p_d = tf_d.paragraphs[0]
    p_d.text = "DATA INGESTION LAYER: Historical Energy Records (CEA, NITI Aayog) + Weather Data (NASA POWER, IMD)"
    p_d.font.name = FONT_FAMILY
    p_d.font.size = Pt(10)
    p_d.font.bold = True
    p_d.font.color.rgb = C_NAVY_BLUE
    p_d.alignment = PP_ALIGN.CENTER

    arr1 = slide11.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.45), Inches(1.98), Inches(0.4), Inches(0.25))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = C_ACCENT_BLUE
    arr1.line.fill.background()

    # 2. Preprocessing & Feature Engineering Layer
    p_box = add_card(slide11, Inches(2.5), Inches(2.25), Inches(8.333), Inches(0.65), bg_color=C_WHITE, border_color=C_BORDER)
    tb_p = slide11.shapes.add_textbox(Inches(2.6), Inches(2.3), Inches(8.133), Inches(0.55))
    tf_p = tb_p.text_frame
    p_p = tf_p.paragraphs[0]
    p_p.text = "PREPROCESSING & FEATURE ENGINEERING: Missing Value Imputation • Date Alignment • Lag (1,2,3,6,12) • Rolling Averages • Cyclical Encoders • MinMax Scaling"
    p_p.font.name = FONT_FAMILY
    p_p.font.size = Pt(9)
    p_p.font.bold = True
    p_p.font.color.rgb = C_TEXT_DARK
    p_p.alignment = PP_ALIGN.CENTER

    arr_d_left = slide11.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.0), Inches(2.93), Inches(0.35), Inches(0.25))
    arr_d_left.fill.solid()
    arr_d_left.fill.fore_color.rgb = C_ACCENT_BLUE
    arr_d_left.line.fill.background()

    arr_s_right = slide11.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(8.9), Inches(2.93), Inches(0.35), Inches(0.25))
    arr_s_right.fill.solid()
    arr_s_right.fill.fore_color.rgb = C_ACCENT_BLUE
    arr_s_right.line.fill.background()

    # 3. Independent Forecasting Modules (Strictly Decoupled)
    # Left Module: Demand Forecasting
    d_mod = add_card(slide11, Inches(1.5), Inches(3.2), Inches(4.8), Inches(1.15), bg_color=RGBColor(240, 249, 255), border_color=RGBColor(186, 230, 253))
    tb_dm = slide11.shapes.add_textbox(Inches(1.6), Inches(3.25), Inches(4.6), Inches(1.05))
    tf_dm = tb_dm.text_frame
    p_dm1 = tf_dm.paragraphs[0]
    p_dm1.text = "MODULE A: ELECTRICITY DEMAND FORECASTING"
    p_dm1.font.name = FONT_FAMILY
    p_dm1.font.size = Pt(10)
    p_dm1.font.bold = True
    p_dm1.font.color.rgb = RGBColor(3, 105, 161)
    
    p_dm2 = tf_dm.add_paragraph()
    p_dm2.text = "• Model: LSTM Sequence Model (Lookback = 3)\n• Outputs: Point Demand (D_hat) & Uncertainty (sigma_D = 391.30 MU)\n• Multi-Level Intervals: 90%, 95%, 99% Demand PIs"
    p_dm2.font.name = FONT_FAMILY
    p_dm2.font.size = Pt(8.5)
    p_dm2.font.color.rgb = C_TEXT_DARK

    # Right Module: Supply Forecasting
    s_mod = add_card(slide11, Inches(7.0), Inches(3.2), Inches(4.8), Inches(1.15), bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))
    tb_sm = slide11.shapes.add_textbox(Inches(7.1), Inches(3.25), Inches(4.6), Inches(1.05))
    tf_sm = tb_sm.text_frame
    p_sm1 = tf_sm.paragraphs[0]
    p_sm1.text = "MODULE B: ELECTRICITY SUPPLY FORECASTING"
    p_sm1.font.name = FONT_FAMILY
    p_sm1.font.size = Pt(10)
    p_sm1.font.bold = True
    p_sm1.font.color.rgb = RGBColor(21, 128, 61)
    
    p_sm2 = tf_sm.add_paragraph()
    p_sm2.text = "• Model: LSTM Sequence Model (8 Sources, Lookback = 6)\n• Outputs: Point Supply (S_hat) & Uncertainty (sigma_S = 996.83 MU)\n• Multi-Level Intervals: 90%, 95%, 99% Supply PIs"
    p_sm2.font.name = FONT_FAMILY
    p_sm2.font.size = Pt(8.5)
    p_sm2.font.color.rgb = C_TEXT_DARK

    # Arrows feeding into Gap Analysis
    arr_g_l = slide11.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.0), Inches(4.38), Inches(0.35), Inches(0.22))
    arr_g_l.fill.solid()
    arr_g_l.fill.fore_color.rgb = C_ACCENT_BLUE
    arr_g_l.line.fill.background()

    arr_g_r = slide11.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(8.9), Inches(4.38), Inches(0.35), Inches(0.22))
    arr_g_r.fill.solid()
    arr_g_r.fill.fore_color.rgb = C_ACCENT_BLUE
    arr_g_r.line.fill.background()

    # 4. Demand–Supply Gap & Uncertainty Engine
    gap_box = add_card(slide11, Inches(1.5), Inches(4.62), Inches(10.3), Inches(0.85), bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER)
    tb_gap = slide11.shapes.add_textbox(Inches(1.6), Inches(4.66), Inches(10.1), Inches(0.75))
    tf_gap = tb_gap.text_frame
    p_g1 = tf_gap.paragraphs[0]
    p_g1.text = "DEMAND–SUPPLY GAP ANALYSIS & PROBABILISTIC UNCERTAINTY PROPAGATION"
    p_g1.font.name = FONT_FAMILY
    p_g1.font.size = Pt(10)
    p_g1.font.bold = True
    p_g1.font.color.rgb = C_NAVY_DARK
    p_g1.alignment = PP_ALIGN.CENTER
    
    p_g2 = tf_gap.add_paragraph()
    p_g2.text = "Gap Equation: G_hat = D_hat - S_hat  |  Combined Uncertainty: sigma_G = sqrt(sigma_D^2 + sigma_S^2) = 1,070.88 MU  |  90%, 95%, 99% Gap Prediction Intervals\nCondition Classification: Shortage (G_hat > 0) • Balanced (G_hat = 0) • Surplus (G_hat < 0)  |  Deterministic Risk: Low (<3000) • Moderate (3000-4500) • High (>4500 MU)"
    p_g2.font.name = FONT_FAMILY
    p_g2.font.size = Pt(8.5)
    p_g2.font.color.rgb = C_TEXT_DARK
    p_g2.alignment = PP_ALIGN.CENTER

    arr_rag = slide11.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.45), Inches(5.5), Inches(0.4), Inches(0.22))
    arr_rag.fill.solid()
    arr_rag.fill.fore_color.rgb = C_ACCENT_BLUE
    arr_rag.line.fill.background()

    # 5. RAG Engine & Output Layer
    rag_box = add_card(slide11, Inches(1.0), Inches(5.75), Inches(11.333), Inches(1.3), bg_color=C_WHITE, border_color=C_BORDER)
    tb_rag = slide11.shapes.add_textbox(Inches(1.1), Inches(5.8), Inches(11.133), Inches(1.2))
    tf_rag = tb_rag.text_frame
    p_r1 = tf_rag.paragraphs[0]
    p_r1.text = "RETRIEVAL-AUGMENTED GENERATION (RAG) & ENERGY PLANNING RECOMMENDATION ENGINE"
    p_r1.font.name = FONT_FAMILY
    p_r1.font.size = Pt(10)
    p_r1.font.bold = True
    p_r1.font.color.rgb = C_ACCENT_TEAL
    
    p_r2 = tf_rag.add_paragraph()
    p_r2.text = "• Knowledge Base: Energy Acts, TNERC Grid Codes, CEA Planning Reports, SLDC Operating Norms, Demand Response Protocols\n• Vector Store & Retrieval: Domain Text Chunking -> Dense Embeddings -> FAISS Vector Database (Top-k Semantic Similarity Search)\n• LLM Synthesis: Forecast Numbers + Risk Severity + Retrieved Regulatory Evidence -> Contextual Energy Planning Actions\n• Final Deliverable: Generation Scheduling, Power Trading, Peaking Plant Dispatch, Maintenance Planning & Grid Reliability Support"
    p_r2.font.name = FONT_FAMILY
    p_r2.font.size = Pt(8.5)
    p_r2.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 12: METHODOLOGY FLOW CHART
    # =========================================================================
    slide12 = add_base_slide(12, "Methodology Flowchart & Stage-by-Stage Pipeline")

    flow_steps = [
        ("Step 1", "Data Collection (CEA, NITI Aayog, NASA POWER, IMD: 129 monthly records)"),
        ("Step 2", "Data Preprocessing & Feature Engineering (Lags 1-12, Rolling Means, Sin/Cos)"),
        ("Step 3", "Independent ML Training: RF, XGBoost, LightGBM, LSTM (80/20 Chronological Split)"),
        ("Step 4", "Model Evaluation & Selection via MAE, RMSE, MAPE, and R² Metrics"),
        ("Step 5", "Point Forecasting for Demand (D_hat) and Source-Wise Generation (S_hat)"),
        ("Step 6", "Residual Uncertainty Quantification: sigma_D = 391.30 MU, sigma_S = 996.83 MU"),
        ("Step 7", "Multi-Level Prediction Intervals Generation: 90%, 95%, and 99% Bounds"),
        ("Step 8", "Demand–Supply Gap & Combined Uncertainty: G_hat = D_hat - S_hat, sigma_G = 1070.88 MU"),
        ("Step 9", "Deterministic Risk Classification (Shortage/Surplus; Low, Moderate, High Impact)"),
        ("Step 10", "RAG Vector Retrieval (FAISS) + LLM Synthesis -> Energy Planning Policy Actions"),
    ]

    for idx, (snum, sdesc) in enumerate(flow_steps):
        top_s = Inches(1.3 + idx * 0.56)
        card_s = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_s, Inches(6.0), Inches(0.48))
        card_s.fill.solid()
        card_s.fill.fore_color.rgb = C_BLUE_LIGHT if idx in [2, 4, 7, 9] else C_WHITE
        card_s.line.color.rgb = C_BLUE_BORDER if idx in [2, 4, 7, 9] else C_BORDER
        card_s.line.width = Pt(1)

        tb_s = slide12.shapes.add_textbox(Inches(0.9), top_s + Inches(0.04), Inches(5.8), Inches(0.4))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = 0
        p_s = tf_s.paragraphs[0]
        p_s.text = f"{snum}: {sdesc}"
        p_s.font.name = FONT_FAMILY
        p_s.font.size = Pt(8.5)
        p_s.font.bold = (idx in [2, 4, 7, 9])
        p_s.font.color.rgb = C_NAVY_DARK

    add_card(slide12, Inches(7.1), Inches(1.3), Inches(5.4), Inches(5.65), title="METHODOLOGICAL STAGES BREAKDOWN", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    tb_mb = slide12.shapes.add_textbox(Inches(7.3), Inches(1.75), Inches(5.0), Inches(5.1))
    tf_mb = tb_mb.text_frame
    tf_mb.word_wrap = True

    stages_desc = [
        ("Phase I: Data Integration & Engineering", 
         "129 monthly rows (Apr 2015 – Dec 2025). Missing values imputed. Features include 3-month and 6-month historical lookbacks, rolling statistics, and cyclical calendar harmonics."),
        
        ("Phase II: Decoupled ML Modeling", 
         "Four candidate architectures evaluated (RF, XGBoost, LightGBM, LSTM). Models trained separately on demand and 8 supply sources. LSTM demonstrated superior long-term dependency capture."),
        
        ("Phase III: Residual Uncertainty & Gap Math", 
         "Residual distributions quantify empirical standard errors (sigma_e). Symmetric prediction intervals generated using standard normal critical multipliers (z=1.645, 1.960, 2.576). Gap error propagated orthogonally."),
        
        ("Phase IV: RAG-Grounded Decision Support", 
         "Forecast outputs and risk levels form a structured semantic query. FAISS retrieves top-k relevant energy policy clauses. LLM generates explainable dispatch and resource planning recommendations.")
    ]

    for st_idx, (st_t, st_d) in enumerate(stages_desc):
        p_st = tf_mb.add_paragraph() if st_idx > 0 else tf_mb.paragraphs[0]
        p_st.text = st_t
        p_st.font.name = FONT_FAMILY
        p_st.font.size = Pt(10)
        p_st.font.bold = True
        p_st.font.color.rgb = C_NAVY_BLUE
        p_st.space_before = Pt(6) if st_idx > 0 else Pt(0)
        
        p_std = tf_mb.add_paragraph()
        p_std.text = st_d
        p_std.font.name = FONT_FAMILY
        p_std.font.size = Pt(9)
        p_std.font.color.rgb = C_TEXT_DARK
        p_std.space_before = Pt(2)

    # =========================================================================
    # SLIDE 13: MATERIALS, METHODS & DATA ANALYSIS
    # =========================================================================
    slide13 = add_base_slide(13, "Materials, Methods & Experimental Data Analysis")

    add_card(slide13, Inches(0.8), Inches(1.3), Inches(5.6), Inches(2.7), title="DATASET SPECIFICATIONS & SOURCES", bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
    t_shape13a = slide13.shapes.add_table(5, 3, Inches(0.95), Inches(1.75), Inches(5.3), Inches(2.1))
    tbl13a = t_shape13a.table
    tbl13a.columns[0].width = Inches(1.8)
    tbl13a.columns[1].width = Inches(1.9)
    tbl13a.columns[2].width = Inches(1.6)
    
    ds_headers = ["Domain / Dataset", "Official Source", "Records / Range"]
    for idx, h in enumerate(ds_headers):
        tbl13a.cell(0, idx).text = h
    
    ds_rows = [
        ("Electricity Demand", "NITI Aayog / CEA", "129 Months (2015-2025)"),
        ("Source-Wise Supply", "NITI Aayog / CEA (8 sources)", "129 Months (2015-2025)"),
        ("Temperature & Humidity", "NASA POWER Portal", "129 Months (Tamil Nadu)"),
        ("Rainfall & Solar Irradiance", "IMD / National Water Portal", "129 Months (Aggregated)")
    ]
    for r_idx, r_val in enumerate(ds_rows):
        for c_idx, c_val in enumerate(r_val):
            tbl13a.cell(r_idx + 1, c_idx).text = c_val
    style_table(tbl13a)

    add_card(slide13, Inches(6.6), Inches(1.3), Inches(5.9), Inches(2.7), title="MATHEMATICAL FORMULATIONS", bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
    tb_math = slide13.shapes.add_textbox(Inches(6.8), Inches(1.75), Inches(5.5), Inches(2.1))
    tf_math = tb_math.text_frame
    tf_math.word_wrap = True
    
    math_eqs = [
        ("Residual Standard Error:", "sigma_e = sqrt( sum( (y_act - y_pred)^2 ) / (N - 1) )"),
        ("Prediction Intervals (PI):", "PI_(1-alpha) = [ y_hat - z_(1-alpha/2) * sigma,  y_hat + z_(1-alpha/2) * sigma ]"),
        ("Critical Z-Multipliers:", "90% -> z=1.645  |  95% -> z=1.960  |  99% -> z=2.576"),
        ("Combined Gap & Sigma:", "G_hat = D_hat - S_hat  |  sigma_Gap = sqrt( sigma_Demand^2 + sigma_Supply^2 )"),
        ("Deterministic Risk Bands:", "Low: Gap < 3000 MU  |  Moderate: 3000-4500 MU  |  High: Gap > 4500 MU")
    ]
    for m_idx, (mlbl, meq) in enumerate(math_eqs):
        p = tf_math.add_paragraph() if m_idx > 0 else tf_math.paragraphs[0]
        p.text = f"•  {mlbl}  "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_DARK
        p.space_before = Pt(3) if m_idx > 0 else Pt(0)
        run = p.add_run()
        run.text = meq
        run.font.name = FONT_FAMILY
        run.font.size = Pt(9)
        run.font.bold = False
        run.font.color.rgb = C_NAVY_BLUE

    add_card(slide13, Inches(0.8), Inches(4.15), Inches(5.6), Inches(2.85), title="DEMAND MODEL PERFORMANCE COMPARISON", bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
    t_shape13b = slide13.shapes.add_table(5, 5, Inches(0.95), Inches(4.6), Inches(5.3), Inches(2.25))
    tbl13b = t_shape13b.table
    tbl13b.columns[0].width = Inches(1.5)
    tbl13b.columns[1].width = Inches(0.95)
    tbl13b.columns[2].width = Inches(0.95)
    tbl13b.columns[3].width = Inches(1.0)
    tbl13b.columns[4].width = Inches(0.9)

    m_headers = ["Model", "MAE (MU)", "RMSE (MU)", "MAPE (%)", "R² Score"]
    for idx, h in enumerate(m_headers):
        tbl13b.cell(0, idx).text = h

    demand_perf = [
        ("Random Forest", "741.63", "909.69", "6.56%", "0.31"),
        ("XGBoost", "619.18", "724.51", "5.56%", "0.34"),
        ("LightGBM", "634.18", "748.24", "5.66%", "0.29"),
        ("LSTM (Selected)", "368.72", "412.93", "3.41%", "0.78")
    ]
    for r_idx, r_val in enumerate(demand_perf):
        for c_idx, c_val in enumerate(r_val):
            tbl13b.cell(r_idx + 1, c_idx).text = c_val
    style_table(tbl13b)
    for c in range(5):
        tbl13b.cell(4, c).fill.solid()
        tbl13b.cell(4, c).fill.fore_color.rgb = RGBColor(220, 252, 231)

    add_card(slide13, Inches(6.6), Inches(4.15), Inches(5.9), Inches(2.85), title="SUPPLY MODEL PERFORMANCE COMPARISON", bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
    t_shape13c = slide13.shapes.add_table(5, 5, Inches(6.75), Inches(4.6), Inches(5.6), Inches(2.25))
    tbl13c = t_shape13c.table
    tbl13c.columns[0].width = Inches(1.6)
    tbl13c.columns[1].width = Inches(1.0)
    tbl13c.columns[2].width = Inches(1.0)
    tbl13c.columns[3].width = Inches(1.05)
    tbl13c.columns[4].width = Inches(0.95)

    for idx, h in enumerate(m_headers):
        tbl13c.cell(0, idx).text = h

    supply_perf = [
        ("Random Forest", "823.30", "1220.90", "6.89%", "0.31"),
        ("XGBoost", "847.76", "1215.30", "7.13%", "0.31"),
        ("LightGBM", "802.74", "1157.40", "6.74%", "0.38"),
        ("LSTM (Selected)", "731.56", "973.93", "6.30%", "0.56")
    ]
    for r_idx, r_val in enumerate(supply_perf):
        for c_idx, c_val in enumerate(r_val):
            tbl13c.cell(r_idx + 1, c_idx).text = c_val
    style_table(tbl13c)
    for c in range(5):
        tbl13c.cell(4, c).fill.solid()
        tbl13c.cell(4, c).fill.fore_color.rgb = RGBColor(220, 252, 231)

    # =========================================================================
    # SLIDE 14: VALIDATION, TIMELINE & BUDGET
    # =========================================================================
    slide14 = add_base_slide(14, "Experimental Validation, Project Timeline & Budget")

    add_card(slide14, Inches(0.8), Inches(1.3), Inches(11.733), Inches(2.45), title="HOLD-OUT FUTURE VALIDATION BENCHMARK (JANUARY – MARCH 2026)", bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
    
    t_shape14a = slide14.shapes.add_table(4, 9, Inches(0.95), Inches(1.7), Inches(11.433), Inches(1.9))
    tbl14a = t_shape14a.table
    tbl14a.columns[0].width = Inches(1.0)
    tbl14a.columns[1].width = Inches(1.2)
    tbl14a.columns[2].width = Inches(1.3)
    tbl14a.columns[3].width = Inches(1.2)
    tbl14a.columns[4].width = Inches(1.3)
    tbl14a.columns[5].width = Inches(1.1)
    tbl14a.columns[6].width = Inches(1.2)
    tbl14a.columns[7].width = Inches(1.8)
    tbl14a.columns[8].width = Inches(1.333)

    val_headers = ["Month", "Actual Dem", "Pred Demand", "Actual Sup", "Pred Supply", "Actual Gap", "Pred Gap", "95% Gap PI [Lower, Upper]", "Condition / Risk"]
    for idx, h in enumerate(val_headers):
        tbl14a.cell(0, idx).text = h

    val_rows = [
        ("Jan 2026", "10,067.00", "11,047.51", "10,189.56", "8,809.63", "-122.56", "+2,237.88", "[138.99,  4,336.77] MU", "Shortage (Low Risk)"),
        ("Feb 2026", "10,125.00", "12,308.25", "10,405.60", "9,229.95", "-280.60", "+3,078.30", "[979.41,  5,177.19] MU", "Shortage (Moderate)"),
        ("Mar 2026", "12,233.00", "12,594.89", "11,247.51", "9,574.64", "+985.49", "+3,020.25", "[921.36,  5,119.14] MU", "Shortage (Moderate)")
    ]

    for r_idx, r_val in enumerate(val_rows):
        for c_idx, c_val in enumerate(r_val):
            tbl14a.cell(r_idx + 1, c_idx).text = c_val
    style_table(tbl14a)

    add_card(slide14, Inches(0.8), Inches(3.9), Inches(6.8), Inches(3.1), title="PROJECT IMPLEMENTATION TIMELINE (PHASE I & II)", bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
    t_shape14b = slide14.shapes.add_table(6, 4, Inches(0.95), Inches(4.35), Inches(6.5), Inches(2.5))
    tbl14b = t_shape14b.table
    tbl14b.columns[0].width = Inches(1.8)
    tbl14b.columns[1].width = Inches(2.5)
    tbl14b.columns[2].width = Inches(1.1)
    tbl14b.columns[3].width = Inches(1.1)

    time_headers = ["Project Phase", "Key Deliverables & Milestones", "Duration", "Status"]
    for idx, h in enumerate(time_headers):
        tbl14b.cell(0, idx).text = h

    timeline_data = [
        ("Phase I: Initiation", "Literature Survey, Dataset Collection & Cleaning", "Months 1–2", "Completed"),
        ("Phase I: Modeling", "Model Training (RF, XGB, LGBM, LSTM) & Tuning", "Months 3–4", "Completed"),
        ("Phase I: Analytics", "Uncertainty Quantification (90/95/99% PIs) & Gap Math", "Month 5", "Completed"),
        ("Phase I: Review", "RAG Pipeline Prototype & IEEE Paper Preparation", "Month 6", "Completed"),
        ("Phase II: Deployment", "Interactive Web UI, Multi-Agent RAG & Live APIs", "Months 7–10", "Planned")
    ]
    for r_idx, r_val in enumerate(timeline_data):
        for c_idx, c_val in enumerate(r_val):
            tbl14b.cell(r_idx + 1, c_idx).text = c_val
    style_table(tbl14b)

    add_card(slide14, Inches(7.8), Inches(3.9), Inches(4.733), Inches(3.1), title="ACADEMIC PROJECT BUDGET", bg_color=C_WHITE, border_color=C_BORDER, title_size=11.5)
    t_shape14c = slide14.shapes.add_table(5, 3, Inches(7.95), Inches(4.35), Inches(4.433), Inches(2.5))
    tbl14c = t_shape14c.table
    tbl14c.columns[0].width = Inches(2.2)
    tbl14c.columns[1].width = Inches(1.233)
    tbl14c.columns[2].width = Inches(1.0)

    b_headers = ["Resource Category", "Provisioning", "Cost (INR)"]
    for idx, h in enumerate(b_headers):
        tbl14c.cell(0, idx).text = h

    budget_data = [
        ("Hardware / Compute", "Institutional Lab", "₹ 0 (In-house)"),
        ("Datasets (CEA/NASA)", "Open Public Portals", "₹ 0 (Open Data)"),
        ("Software Frameworks", "Open-Source Python", "₹ 0 (FOSS)"),
        ("Vector DB & Local LLM", "FAISS / Ollama Local", "₹ 0 (Zero-cost)")
    ]
    for r_idx, r_val in enumerate(budget_data):
        for c_idx, c_val in enumerate(r_val):
            tbl14c.cell(r_idx + 1, c_idx).text = c_val
    style_table(tbl14c)

    # =========================================================================
    # SLIDE 15: PHASE II WORK & FUTURE DIRECTIONS
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    bg15 = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg15.fill.solid()
    bg15.fill.fore_color.rgb = C_SLATE_BG
    bg15.line.fill.background()

    top_bar15 = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.12))
    top_bar15.fill.solid()
    top_bar15.fill.fore_color.rgb = C_NAVY_DARK
    top_bar15.line.color.rgb = C_NAVY_DARK

    cat_box15 = slide15.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.0), Inches(0.28))
    tf_c15 = cat_box15.text_frame
    p_c15 = tf_c15.paragraphs[0]
    p_c15.text = "FINAL YEAR PROJECT REVIEW — PHASE II & FUTURE SCOPE"
    p_c15.font.name = FONT_FAMILY
    p_c15.font.size = Pt(9.5)
    p_c15.font.bold = True
    p_c15.font.color.rgb = RGBColor(147, 197, 253)

    title_box15 = slide15.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.0), Inches(0.65))
    tf_t15 = title_box15.text_frame
    p_t15 = tf_t15.paragraphs[0]
    p_t15.text = "Phase II Implementation Plan, Future Scope & Conclusions"
    p_t15.font.name = FONT_FAMILY
    p_t15.font.size = Pt(20)
    p_t15.font.bold = True
    p_t15.font.color.rgb = C_WHITE

    accent_line15 = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.12), Inches(13.333), Inches(0.04))
    accent_line15.fill.solid()
    accent_line15.fill.fore_color.rgb = C_ACCENT_BLUE
    accent_line15.line.color.rgb = C_ACCENT_BLUE

    add_card(slide15, Inches(0.8), Inches(1.35), Inches(5.75), Inches(3.7), title="PLANNED PHASE II IMPLEMENTATION WORK", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    tb_p2 = slide15.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.35), Inches(3.05))
    tf_p2 = tb_p2.text_frame
    tf_p2.word_wrap = True

    p2_pts = [
        ("Interactive Web Dashboard:", "Deploy a production-grade Streamlit/FastAPI web interface featuring live scenario sliders, interval toggles, and condition gauges."),
        ("Real-Time Data Ingestion:", "Integrate automated live weather APIs (OpenWeather / IMD) and real-time SCADA generation data streams."),
        ("Multi-Agent RAG Architecture:", "Implement multi-agent cooperative verification where specialized LLM agents debate and refine dispatch recommendations."),
        ("Hyperparameter Optimization:", "Explore Bayesian optimization for deep recurrent networks and automated residual calibration across shifting seasons.")
    ]
    for idx, (p_title, p_desc) in enumerate(p2_pts):
        p = tf_p2.add_paragraph() if idx > 0 else tf_p2.paragraphs[0]
        p.text = f"•  {p_title} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_BLUE
        p.space_before = Pt(4) if idx > 0 else Pt(0)
        run = p.add_run()
        run.text = p_desc
        run.font.name = FONT_FAMILY
        run.font.size = Pt(9)
        run.font.bold = False
        run.font.color.rgb = C_TEXT_DARK

    add_card(slide15, Inches(6.8), Inches(1.35), Inches(5.75), Inches(3.7), title="FUTURE RESEARCH DIRECTIONS", bg_color=C_WHITE, border_color=C_BORDER, title_size=12)
    tb_fs = slide15.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(5.35), Inches(3.05))
    tf_fs = tb_fs.text_frame
    tf_fs.word_wrap = True

    fs_pts = [
        ("Advanced Transformer Architectures:", "Investigate Temporal Fusion Transformers (TFT) and Informer models for high-resolution sub-hourly load forecasting."),
        ("Dynamic Conformal Prediction:", "Implement non-parametric conformal prediction intervals to capture extreme climate anomalies without assuming normal residuals."),
        ("Battery Energy Storage (BESS) Dispatch:", "Incorporate automated optimal BESS charge/discharge scheduling algorithms directly into the RAG recommendation loop."),
        ("Inter-State Grid Power Wheeling:", "Extend gap analytics to regional inter-state exchange corridors and dynamic spot market pricing models.")
    ]
    for idx, (f_title, f_desc) in enumerate(fs_pts):
        p = tf_fs.add_paragraph() if idx > 0 else tf_fs.paragraphs[0]
        p.text = f"•  {f_title} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT_TEAL
        p.space_before = Pt(4) if idx > 0 else Pt(0)
        run = p.add_run()
        run.text = f_desc
        run.font.name = FONT_FAMILY
        run.font.size = Pt(9)
        run.font.bold = False
        run.font.color.rgb = C_TEXT_DARK

    thank_card = add_card(slide15, Inches(0.8), Inches(5.25), Inches(11.75), Inches(1.75), bg_color=C_NAVY_DARK, border_color=C_NAVY_DARK)
    tb_ty = slide15.shapes.add_textbox(Inches(1.0), Inches(5.35), Inches(11.35), Inches(1.55))
    tf_ty = tb_ty.text_frame
    tf_ty.word_wrap = True
    
    p_ty1 = tf_ty.paragraphs[0]
    p_ty1.text = "THANK YOU!"
    p_ty1.font.name = FONT_FAMILY
    p_ty1.font.size = Pt(22)
    p_ty1.font.bold = True
    p_ty1.font.color.rgb = RGBColor(147, 197, 253)
    p_ty1.alignment = PP_ALIGN.CENTER

    p_ty2 = tf_ty.add_paragraph()
    p_ty2.text = "Questions & Discussion"
    p_ty2.font.name = FONT_FAMILY
    p_ty2.font.size = Pt(14)
    p_ty2.font.bold = True
    p_ty2.font.color.rgb = C_WHITE
    p_ty2.alignment = PP_ALIGN.CENTER
    p_ty2.space_before = Pt(2)

    p_ty3 = tf_ty.add_paragraph()
    p_ty3.text = "Confidence-Driven Probabilistic Machine Learning Framework for Electricity Demand–Supply Gap Forecasting and RAG-Based Energy Planning\nDepartment of Information Technology  •  Academic Year 2026–2027"
    p_ty3.font.name = FONT_FAMILY
    p_ty3.font.size = Pt(10)
    p_ty3.font.color.rgb = RGBColor(203, 213, 225)
    p_ty3.alignment = PP_ALIGN.CENTER
    p_ty3.space_before = Pt(4)

    # Save to multiple locations for easy access
    out_paths = [
        r'd:\Final Year Project 2027\Electricity_Forecasting_Final_Review.pptx',
        r'd:\Final Year Project 2027\Electricity Demand and Supply Forecasting\Electricity_Forecasting_Final_Review.pptx'
    ]
    for p in out_paths:
        prs.save(p)
        print(f"Presentation successfully saved to: {p}")

if __name__ == '__main__':
    build_presentation()
